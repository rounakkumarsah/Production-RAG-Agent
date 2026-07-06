import fitz
doc = fitz.open()
page = doc.new_page()
page.insert_text((50,50), 'Test PDF')
doc.save('/tmp/test.pdf')
doc.close()

from docx import Document
d = Document()
d.add_paragraph('Test DOCX')
d.save('/tmp/test.docx')

from pptx import Presentation
p = Presentation()
slide = p.slides.add_slide(p.slide_layouts[0])
slide.shapes.title.text = 'Test PPTX'
p.save('/tmp/test.pptx')

from openpyxl import Workbook
wb = Workbook()
ws = wb.active
ws['A1'] = 'Test XLSX'
wb.save('/tmp/test.xlsx')

with open('/tmp/test.csv', 'w') as f:
    f.write('col1,col2\nTest,CSV\n')

with open('/tmp/test.txt', 'w') as f:
    f.write('Test TXT\n')

with open('/tmp/test.md', 'w') as f:
    f.write('# Test MD\n')
