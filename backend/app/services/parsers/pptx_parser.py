import io
from pptx import Presentation
from app.services.parsers.base import DocumentParser, ExtractionResult, ExtractedPage
from app.utils.logging import get_logger

logger = get_logger(__name__)

class PptxParser(DocumentParser):
    def parse(self, content: bytes, filename: str) -> ExtractionResult:
        try:
            prs = Presentation(io.BytesIO(content))
        except Exception as exc:
            raise ValueError(f"Cannot open PPTX '{filename}': {exc}") from exc

        pages: list[ExtractedPage] = []
        cursor = 0
        total_pages = len(prs.slides)
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            
            page_text = "\n".join(slide_text).strip()
            if not page_text:
                continue

            start = cursor
            end = start + len(page_text)
            
            pages.append(ExtractedPage(
                page_number=i + 1,
                text=page_text,
                char_start=start,
                char_end=end
            ))
            cursor = end + 2  # account for \n\n

        if not pages:
            raise ValueError(f"PPTX '{filename}' contains no extractable text.")
            
        full_text = "\n\n".join(p.text for p in pages)
        
        return ExtractionResult(
            pages=pages,
            full_text=full_text,
            page_count=total_pages,
            char_count=len(full_text),
            file_type="pptx",
            parser_used="python-pptx",
            ocr_used=False,
            ocr_engine=None,
            extraction_method="native"
        )
