import os
import requests
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
import pandas as pd
from fpdf import FPDF

API_URL = "http://localhost:8000"

def create_files():
    os.makedirs("test_files", exist_ok=True)
    
    # 1. TXT
    with open("test_files/test.txt", "w") as f:
        f.write("This is a simple text file for testing TXT parsing.")

    # 2. Markdown
    with open("test_files/test.md", "w") as f:
        f.write("# Markdown Test\n\nThis is a markdown file with a **bold** word.")

    # 3. CSV
    df = pd.DataFrame({"Name": ["Alice", "Bob"], "Age": [25, 30]})
    df.to_csv("test_files/test.csv", index=False)

    # 4. DOCX
    doc = Document()
    doc.add_paragraph("This is a Word document for testing DOCX parsing.")
    doc.save("test_files/test.docx")

    # 5. PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PPTX Test"
    slide.placeholders[1].text = "This is a PowerPoint presentation."
    prs.save("test_files/test.pptx")

    # 6. XLSX
    wb = Workbook()
    ws = wb.active
    ws.append(["Item", "Cost"])
    ws.append(["Apple", 1.2])
    wb.save("test_files/test.xlsx")

    # 7. PDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="This is a standard PDF file for testing PDF parsing.", ln=1)
    pdf.output("test_files/test.pdf")
    
def test_uploads():
    results = {}
    for filename in os.listdir("test_files"):
        filepath = os.path.join("test_files", filename)
        with open(filepath, "rb") as f:
            resp = requests.post(f"{API_URL}/upload", files={"files": (filename, f)})
            print(f"Uploaded {filename}: {resp.status_code}")
            try:
                data = resp.json()
                if "documents" in data and len(data["documents"]) > 0:
                    doc = data["documents"][0]
                    results[filename] = {"status": doc["status"], "chunks": doc.get("chunk_count", 0)}
                else:
                    results[filename] = {"error": "Invalid response"}
            except Exception as e:
                results[filename] = {"error": str(e)}
    print(results)

if __name__ == "__main__":
    create_files()
    test_uploads()
